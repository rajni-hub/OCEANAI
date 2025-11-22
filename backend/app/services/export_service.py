"""
Export service - Business logic for document export
"""
from sqlalchemy.orm import Session
from fastapi import HTTPException, status
from typing import Dict, Any, Optional
from io import BytesIO
from datetime import datetime
import logging

from app.models.project import Project, DocumentType
from app.models.user import User
from app.services.document_service import get_document

# Document generation libraries
from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import PP_PLACEHOLDER


def export_word_document(
    db: Session,
    project: Project,
    user: User
) -> BytesIO:
    """
    Export Word document (.docx) with latest refined content
    
    Args:
        db: Database session
        project: Project object
        user: Authenticated user
        
    Returns:
        BytesIO object containing the .docx file
        
    Raises:
        HTTPException: If validation fails or export fails
    """
    if project.document_type != DocumentType.WORD:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Word export is only available for Word documents"
        )
    
    # Verify project belongs to user
    # Handle SQLite string IDs vs PostgreSQL UUID objects
    from app.core.config import settings
    if 'sqlite' in settings.DATABASE_URL.lower():
        if str(project.user_id) != str(user.id):
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Project does not belong to user"
            )
    else:
        if project.user_id != user.id:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Project does not belong to user"
            )
    
    # Get document
    document = get_document(db=db, project=project)
    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found. Please configure and generate content first."
        )
    
    # Get structure and content
    structure = document.structure
    content = document.content or {}
    
    if not structure or "sections" not in structure:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document structure is invalid"
        )
    
    sections = structure.get("sections", [])
    if not sections:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document has no sections"
        )
    
    # Sort sections by order
    sections = sorted(sections, key=lambda x: x.get("order", 0))
    
    try:
        # Create Word document
        doc = DocxDocument()
        
        # Set document properties
        core_props = doc.core_properties
        core_props.title = project.title
        core_props.author = user.email
        core_props.comments = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add title
        title_paragraph = doc.add_heading(project.title, level=1)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add main topic subtitle
        if project.main_topic:
            topic_paragraph = doc.add_paragraph(project.main_topic)
            topic_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_paragraph_format = topic_paragraph.runs[0].font
            topic_paragraph_format.size = Pt(12)
            topic_paragraph_format.italic = True
        
        # Add spacing
        doc.add_paragraph()
        
        # Process each section
        for section in sections:
            section_id = section.get("id")
            section_title = section.get("title", "Untitled Section")
            
            # Add section heading
            heading = doc.add_heading(section_title, level=2)
            
            # Get section content
            section_content = content.get(section_id, "")
            
            if section_content:
                # Split content into paragraphs (by double newlines or single newlines)
                paragraphs = section_content.split("\n\n")
                
                for para_text in paragraphs:
                    if para_text.strip():
                        # Check if it's a bullet point list
                        lines = para_text.split("\n")
                        is_bullet_list = any(
                            line.strip().startswith(("-", "•", "*", "1.", "2.", "3."))
                            for line in lines
                            if line.strip()
                        )
                        
                        if is_bullet_list:
                            # Add as bullet list
                            for line in lines:
                                line = line.strip()
                                if line:
                                    # Remove bullet markers
                                    for marker in ["-", "•", "*"]:
                                        if line.startswith(marker):
                                            line = line[1:].strip()
                                            break
                                    
                                    # Remove numbered list markers
                                    if line and line[0].isdigit() and "." in line[:3]:
                                        parts = line.split(".", 1)
                                        if len(parts) > 1:
                                            line = parts[1].strip()
                                    
                                    if line:
                                        para = doc.add_paragraph(line, style='List Bullet')
                        else:
                            # Add as regular paragraph
                            para = doc.add_paragraph(para_text.strip())
                            para_format = para.paragraph_format
                            para_format.space_after = Pt(6)
            else:
                # No content available
                para = doc.add_paragraph("[Content not generated]")
                para_format = para.runs[0].font
                para_format.italic = True
                para_format.color = DocxRGBColor(128, 128, 128)
            
            # Add spacing between sections
            doc.add_paragraph()
        
        # Save to BytesIO
        file_stream = BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return file_stream
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error exporting Word document: {str(e)}"
        )


def export_powerpoint_document(
    db: Session,
    project: Project,
    user: User
) -> BytesIO:
    """
    Export PowerPoint document (.pptx) with latest refined content
    
    Args:
        db: Database session
        project: Project object
        user: Authenticated user
        
    Returns:
        BytesIO object containing the .pptx file
        
    Raises:
        HTTPException: If validation fails or export fails
    """
    if project.document_type != DocumentType.POWERPOINT:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="PowerPoint export is only available for PowerPoint documents"
        )
    
    # Verify project belongs to user
    # Handle SQLite string IDs vs PostgreSQL UUID objects
    from app.core.config import settings
    if 'sqlite' in settings.DATABASE_URL.lower():
        if str(project.user_id) != str(user.id):
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Project does not belong to user"
            )
    else:
        if project.user_id != user.id:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Project does not belong to user"
            )
    
    # Get document
    document = get_document(db=db, project=project)
    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found. Please configure and generate content first."
        )
    
    # Get structure and content
    structure = document.structure
    content = document.content or {}
    
    if not structure or "slides" not in structure:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document structure is invalid"
        )
    
    slides = structure.get("slides", [])
    if not slides:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Document has no slides"
        )
    
    # Sort slides by order
    slides = sorted(slides, key=lambda x: x.get("order", 0))
    
    try:
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set presentation properties
        prs.core_properties.title = project.title
        prs.core_properties.author = user.email
        prs.core_properties.comments = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Process each slide
        for slide_data in slides:
            slide_id = slide_data.get("id")
            slide_title = slide_data.get("title", "Untitled Slide")
            
            # Get slide content
            slide_content = content.get(slide_id, "")
            
            # Create slide with title and content layout
            # Layout 1 is "Title and Content" which has both title and content placeholders
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
            
            # Set content - find content placeholder (not title)
            content_shape = None
            for placeholder in slide.placeholders:
                ph_type = placeholder.placeholder_format.type
                # Any placeholder that's not a title is a content placeholder
                if ph_type != PP_PLACEHOLDER.TITLE:
                    content_shape = placeholder
                    break
            
            # Fallback: try to get placeholder at index 1 if available and not title
            if not content_shape and len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                if ph.placeholder_format.type != PP_PLACEHOLDER.TITLE:
                    content_shape = ph
            
            # If still no content shape, use a text box
            if not content_shape:
                # Add a text box for content
                left = PptPt(0.5 * 914400)  # 0.5 inches in EMU
                top = PptPt(1.5 * 914400)   # 1.5 inches in EMU
                width = prs.slide_width - (left * 2)
                height = prs.slide_height - top - PptPt(0.5 * 914400)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
            else:
                text_frame = content_shape.text_frame
            
            text_frame.word_wrap = True
            
            if slide_content:
                # First, check if the entire content is a bullet list
                all_lines = slide_content.split("\n")
                all_bullet_lines = [
                    line.strip() for line in all_lines 
                    if line.strip() and line.strip().startswith(("-", "•", "*", "1.", "2.", "3."))
                ]
                
                # If we have bullet points, process them all together
                if all_bullet_lines:
                    # Clear existing paragraphs first
                    text_frame.clear()
                    
                    # Add all bullet points
                    for idx, line in enumerate(all_bullet_lines):
                        # Remove bullet markers
                        clean_line = line
                        for marker in ["-", "•", "*"]:
                            if clean_line.startswith(marker):
                                clean_line = clean_line[1:].strip()
                                break
                        
                        # Remove numbered list markers
                        if clean_line and clean_line[0].isdigit() and "." in clean_line[:3]:
                            parts = clean_line.split(".", 1)
                            if len(parts) > 1:
                                clean_line = parts[1].strip()
                        
                        if clean_line:
                            p = text_frame.add_paragraph()
                            p.text = clean_line
                            p.level = 0
                            p.font.size = PptPt(18)
                else:
                    # Not a bullet list, process as regular paragraphs
                    paragraphs = slide_content.split("\n\n")
                    
                    for i, para_text in enumerate(paragraphs):
                        if para_text.strip():
                            if i == 0 and len(text_frame.paragraphs) == 1:
                                # First paragraph already exists
                                p = text_frame.paragraphs[0]
                                p.text = para_text.strip()
                                p.font.size = PptPt(18)
                            else:
                                p = text_frame.add_paragraph()
                                p.text = para_text.strip()
                                p.font.size = PptPt(18)
            else:
                # No content available
                if len(text_frame.paragraphs) == 1:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = "[Content not generated]"
                p.font.italic = True
                p.font.size = PptPt(14)
                p.font.color.rgb = PptRGBColor(128, 128, 128)
        
        # Save to BytesIO
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        logger = logging.getLogger(__name__)
        logger.error(f"Error exporting PowerPoint document: {str(e)}\n{error_trace}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error exporting PowerPoint document: {str(e)}"
        )


def get_export_filename(project: Project) -> str:
    """
    Generate export filename based on project
    
    Args:
        project: Project object
        
    Returns:
        Filename string
    """
    # Sanitize project title for filename
    safe_title = "".join(c for c in project.title if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_title = safe_title.replace(' ', '_')
    
    if project.document_type == DocumentType.WORD:
        extension = "docx"
    else:
        extension = "pptx"
    
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    return f"{safe_title}_{timestamp}.{extension}"
