"""
Template applicator - Apply template styles to exported documents
"""
from typing import Dict, Any, Optional
from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptPt, Inches
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from app.models.template import Template
from app.services.template_service import get_template_config, hex_to_rgb


def apply_template_to_word(
    doc: DocxDocument,
    template: Template,
    project_title: str,
    main_topic: Optional[str] = None
):
    """
    Apply template styles to Word document
    
    Args:
        doc: python-docx Document object
        template: Template object
        project_title: Project title text
        main_topic: Optional main topic text
    """
    config = get_template_config(template)
    
    # Get color palette
    colors = config.get("color_palette", {})
    heading_color = hex_to_rgb(colors.get("heading", colors.get("primary", "#000000")))
    body_color = hex_to_rgb(colors.get("body", colors.get("text", "#000000")))
    
    # Get typography
    typo = config.get("typography", {})
    heading_font = typo.get("heading_font", "Arial")
    body_font = typo.get("body_font", "Calibri")
    heading_size = typo.get("heading_size", 44)
    body_size = typo.get("body_size", 18)
    heading_weight = typo.get("heading_weight", "bold")
    
    # Get spacing
    spacing = config.get("spacing", {})
    section_margin = spacing.get("section_margin", 24)
    paragraph_spacing = spacing.get("paragraph_spacing", 12)
    title_margin_bottom = spacing.get("title_margin_bottom", 18)
    
    # Get styles
    styles_config = config.get("styles", {})
    title_alignment = styles_config.get("title_alignment", "center")
    heading_alignment = styles_config.get("heading_alignment", "left")
    body_alignment = styles_config.get("body_alignment", "left")
    
    # Set document margins if specified
    layout = config.get("layout", {})
    margins = layout.get("document_margins")
    if margins:
        sections = doc.sections
        for section in sections:
            section.top_margin = Inches(margins.get("top", 1))
            section.bottom_margin = Inches(margins.get("bottom", 1))
            section.left_margin = Inches(margins.get("left", 1))
            section.right_margin = Inches(margins.get("right", 1))
    
    # Apply styles to all paragraphs
    for paragraph in doc.paragraphs:
        # Check if it's a heading (has style starting with 'Heading')
        if paragraph.style.name.startswith('Heading'):
            # Heading styles
            for run in paragraph.runs:
                run.font.name = heading_font
                run.font.size = Pt(heading_size)
                run.font.color.rgb = DocxRGBColor(*heading_color)
                if heading_weight == "bold":
                    run.font.bold = True
                
            # Set alignment
            if heading_alignment == "center":
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif heading_alignment == "right":
                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            else:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
            # Set spacing after heading
            paragraph_format = paragraph.paragraph_format
            paragraph_format.space_after = Pt(section_margin)
            
        else:
            # Body text styles
            for run in paragraph.runs:
                run.font.name = body_font
                run.font.size = Pt(body_size)
                run.font.color.rgb = DocxRGBColor(*body_color)
                if typo.get("body_weight") == "bold":
                    run.font.bold = True
            
            # Set alignment
            if body_alignment == "center":
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif body_alignment == "right":
                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            else:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            # Set spacing
            paragraph_format = paragraph.paragraph_format
            paragraph_format.space_after = Pt(paragraph_spacing)


def apply_template_to_powerpoint(
    prs: Presentation,
    template: Template,
    project_title: str
):
    """
    Apply template styles to PowerPoint presentation
    
    Args:
        prs: python-pptx Presentation object
        template: Template object
        project_title: Project title text
    """
    config = get_template_config(template)
    
    # Get color palette
    colors = config.get("color_palette", {})
    heading_color = hex_to_rgb(colors.get("heading", colors.get("primary", "#000000")))
    body_color = hex_to_rgb(colors.get("body", colors.get("text", "#000000")))
    background_color = hex_to_rgb(colors.get("background", "#FFFFFF"))
    
    # Get typography
    typo = config.get("typography", {})
    heading_font = typo.get("heading_font", "Arial")
    body_font = typo.get("body_font", "Calibri")
    heading_size = typo.get("heading_size", 44)
    body_size = typo.get("body_size", 18)
    
    # Get layout
    layout = config.get("layout", {})
    slide_width = layout.get("slide_width", 10)
    slide_height = layout.get("slide_height", 7.5)
    
    # Set slide dimensions
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    # Get styles
    styles_config = config.get("styles", {})
    title_alignment = styles_config.get("title_alignment", "center")
    heading_alignment = styles_config.get("heading_alignment", "left")
    body_alignment = styles_config.get("body_alignment", "left")
    
    # Apply styles to all slides
    for slide in prs.slides:
        # Style title if exists
        if slide.shapes.title:
            title_shape = slide.shapes.title
            if title_shape.has_text_frame:
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = heading_font
                        run.font.size = PptPt(heading_size)
                        run.font.color.rgb = PptRGBColor(*heading_color)
                        run.font.bold = True
                    
                    # Set alignment
                    if title_alignment == "center":
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif title_alignment == "right":
                        paragraph.alignment = PP_ALIGN.RIGHT
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT
                    
                    # Set vertical anchor for title
                    if hasattr(title_shape, 'text_frame'):
                        title_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Style content placeholders
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Skip title shape (already styled)
                if shape == slide.shapes.title:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Check if it's a first-level bullet (likely heading)
                        if paragraph.level == 0 and len(paragraph.runs) > 0:
                            # Could be a heading if it's the first paragraph
                            # For now, apply body styles to all content
                            run.font.name = body_font
                            run.font.size = PptPt(body_size)
                            run.font.color.rgb = PptRGBColor(*body_color)
                        else:
                            run.font.name = body_font
                            run.font.size = PptPt(body_size)
                            run.font.color.rgb = PptRGBColor(*body_color)
                    
                    # Set alignment
                    if body_alignment == "center":
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif body_alignment == "right":
                        paragraph.alignment = PP_ALIGN.RIGHT
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT

